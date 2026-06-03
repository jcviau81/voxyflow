"""
Extensible document parser system for Voxyflow RAG pipeline.

Phase 1: .txt, .md, .markdown support.
Phase 2: PDF, DOCX, XLSX/CSV support (graceful — skipped if deps missing).
Phase 3: PPTX (python-pptx) + HTML (bs4/regex) parsers, plus a `best_effort_text`
         decoder for arbitrary text-like files. Images and scanned PDFs are
         handled out-of-band via the vision pipeline (see document_ingest).
"""

import re
import logging
from abc import ABC, abstractmethod
from dataclasses import dataclass, field
from pathlib import Path

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Exceptions
# ---------------------------------------------------------------------------


class UnsupportedFileType(Exception):
    """Raised when no parser is registered for the given file extension."""

    def __init__(self, extension: str):
        self.extension = extension
        super().__init__(
            f"No dedicated parser for file type: {extension!r}. "
            "Dedicated formats: .txt, .md, .pdf, .docx, .xlsx, .csv, .pptx, .html. "
            "Other types fall back to best-effort text extraction or are stored as-is."
        )


# ---------------------------------------------------------------------------
# Data model
# ---------------------------------------------------------------------------


@dataclass
class ParsedDocument:
    """Result of parsing a document file."""

    text: str
    """Full extracted plain text."""

    chunks: list[str] = field(default_factory=list)
    """Text split into overlapping chunks (~500 chars, 50 char overlap)."""

    metadata: dict = field(default_factory=dict)
    """filename, filetype, size_bytes, chunk_count, etc."""


# ---------------------------------------------------------------------------
# Chunking helpers
# ---------------------------------------------------------------------------

_CHUNK_MIN = 100
_CHUNK_MAX = 600
_CHUNK_OVERLAP = 50

# Simple sentence boundary: end of sentence punctuation followed by whitespace or end
_SENTENCE_END_RE = re.compile(r'(?<=[.!?])\s+')


def _split_into_chunks(text: str) -> list[str]:
    """
    Split text into overlapping chunks.

    Strategy:
    1. Split on paragraph breaks (\\n\\n).
    2. Re-split paragraphs > _CHUNK_MAX chars on sentence boundaries.
    3. Keep chunks between _CHUNK_MIN and _CHUNK_MAX chars.
    4. Add _CHUNK_OVERLAP char overlap between consecutive chunks.
    """
    # Step 1: paragraph split
    raw_paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]

    # Step 2: re-split long paragraphs
    segments: list[str] = []
    for para in raw_paragraphs:
        if len(para) <= _CHUNK_MAX:
            segments.append(para)
        else:
            # Split on sentence boundaries
            sentences = _SENTENCE_END_RE.split(para)
            current = ""
            for sentence in sentences:
                sentence = sentence.strip()
                if not sentence:
                    continue
                if len(current) + len(sentence) + 1 <= _CHUNK_MAX:
                    current = (current + " " + sentence).strip() if current else sentence
                else:
                    if current:
                        segments.append(current)
                    # Sentence itself may be too long — hard cut
                    while len(sentence) > _CHUNK_MAX:
                        segments.append(sentence[:_CHUNK_MAX])
                        sentence = sentence[_CHUNK_MAX - _CHUNK_OVERLAP:]
                    current = sentence
            if current:
                segments.append(current)

    # Step 3: filter by min length
    segments = [s for s in segments if len(s) >= _CHUNK_MIN]

    if not segments:
        # Fallback: if text is very short, return as a single chunk (if non-empty)
        stripped = text.strip()
        return [stripped] if stripped else []

    # Step 4: add overlap
    chunks: list[str] = []
    for i, seg in enumerate(segments):
        if i == 0:
            chunks.append(seg)
        else:
            # Prepend tail of previous segment for context overlap
            prev_tail = segments[i - 1][-_CHUNK_OVERLAP:].strip()
            chunk = (prev_tail + " " + seg).strip() if prev_tail else seg
            # Trim to max if overlap pushed it over
            if len(chunk) > _CHUNK_MAX:
                chunk = chunk[:_CHUNK_MAX]
            chunks.append(chunk)

    return chunks


# ---------------------------------------------------------------------------
# Base parser
# ---------------------------------------------------------------------------


class BaseParser(ABC):
    """Abstract base class for document parsers."""

    supported_extensions: list[str] = []

    @abstractmethod
    def parse(self, content: bytes, filename: str) -> ParsedDocument:
        """Parse raw file bytes into a ParsedDocument."""
        ...


# ---------------------------------------------------------------------------
# Text / Markdown parser
# ---------------------------------------------------------------------------


class TextParser(BaseParser):
    """Parser for plain text and Markdown files."""

    supported_extensions = ['.txt', '.md', '.markdown']

    def parse(self, content: bytes, filename: str) -> ParsedDocument:
        try:
            text = content.decode('utf-8')
        except UnicodeDecodeError:
            text = content.decode('latin-1', errors='replace')

        ext = Path(filename).suffix.lower()
        chunks = _split_into_chunks(text)

        return ParsedDocument(
            text=text,
            chunks=chunks,
            metadata={
                'filename': filename,
                'filetype': ext,
                'size_bytes': len(content),
                'chunk_count': len(chunks),
            },
        )


# ---------------------------------------------------------------------------
# PDF parser (Phase 2)
# ---------------------------------------------------------------------------


class PDFParser(BaseParser):
    """Parser for PDF files. Requires pypdf>=4.0.0."""

    supported_extensions = ['.pdf']

    def parse(self, content: bytes, filename: str) -> ParsedDocument:
        import io
        from pypdf import PdfReader  # type: ignore[import]

        reader = PdfReader(io.BytesIO(content))
        page_texts = []
        for page in reader.pages:
            page_text = page.extract_text() or ""
            page_texts.append(page_text)

        text = "\n\n".join(page_texts)
        chunks = _split_into_chunks(text)

        return ParsedDocument(
            text=text,
            chunks=chunks,
            metadata={
                'filename': filename,
                'filetype': '.pdf',
                'size_bytes': len(content),
                'chunk_count': len(chunks),
                'page_count': len(reader.pages),
            },
        )


# ---------------------------------------------------------------------------
# DOCX parser (Phase 2)
# ---------------------------------------------------------------------------


class DocxParser(BaseParser):
    """Parser for Word documents. Requires python-docx>=1.1.0.

    Only the modern Office Open XML ``.docx`` is supported — python-docx cannot
    read the legacy binary ``.doc`` (OLE) format, so ``.doc`` is intentionally
    not registered here and falls through to best-effort handling in ingest.
    """

    supported_extensions = ['.docx']

    def parse(self, content: bytes, filename: str) -> ParsedDocument:
        import io
        from docx import Document as DocxDocument  # type: ignore[import]

        doc = DocxDocument(io.BytesIO(content))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        text = "\n\n".join(paragraphs)
        chunks = _split_into_chunks(text)

        return ParsedDocument(
            text=text,
            chunks=chunks,
            metadata={
                'filename': filename,
                'filetype': Path(filename).suffix.lower(),
                'size_bytes': len(content),
                'chunk_count': len(chunks),
            },
        )


# ---------------------------------------------------------------------------
# XLSX / CSV parser (Phase 2)
# ---------------------------------------------------------------------------


class XlsxParser(BaseParser):
    """
    Parser for spreadsheets.

    - .xlsx: requires openpyxl>=3.1.0
    - .csv: no external dependency required

    Legacy binary ``.xls`` is not supported (openpyxl reads only the modern XML
    format); it falls through to best-effort handling in ingest.
    """

    supported_extensions = ['.xlsx', '.csv']

    def parse(self, content: bytes, filename: str) -> ParsedDocument:
        ext = Path(filename).suffix.lower()

        if ext == '.csv':
            text = self._parse_csv(content)
        else:
            text = self._parse_xlsx(content)

        chunks = _split_into_chunks(text)

        return ParsedDocument(
            text=text,
            chunks=chunks,
            metadata={
                'filename': filename,
                'filetype': ext,
                'size_bytes': len(content),
                'chunk_count': len(chunks),
            },
        )

    def _parse_csv(self, content: bytes) -> str:
        import csv
        import io

        try:
            text = content.decode('utf-8')
        except UnicodeDecodeError:
            text = content.decode('latin-1', errors='replace')

        reader = csv.reader(io.StringIO(text))
        rows = ["\t".join(row) for row in reader if any(cell.strip() for cell in row)]
        return "\n".join(rows)

    def _parse_xlsx(self, content: bytes) -> str:
        import io
        import openpyxl  # type: ignore[import]

        wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
        sheet_texts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                if any(cell is not None for cell in row):
                    rows.append("\t".join("" if cell is None else str(cell) for cell in row))
            if rows:
                sheet_texts.append(f"[{sheet_name}]\n" + "\n".join(rows))
        wb.close()
        return "\n\n".join(sheet_texts)


# ---------------------------------------------------------------------------
# PPTX parser
# ---------------------------------------------------------------------------


class PptxParser(BaseParser):
    """Parser for PowerPoint decks. Requires python-pptx>=0.6."""

    supported_extensions = ['.pptx']

    def parse(self, content: bytes, filename: str) -> ParsedDocument:
        import io
        from pptx import Presentation  # type: ignore[import]

        prs = Presentation(io.BytesIO(content))
        slide_texts: list[str] = []
        for idx, slide in enumerate(prs.slides, start=1):
            parts: list[str] = []
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    txt = shape.text_frame.text
                    if txt and txt.strip():
                        parts.append(txt.strip())
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        cells = [c.text for c in row.cells]
                        if any(c.strip() for c in cells):
                            parts.append("\t".join(cells))
            if getattr(slide, "has_notes_slide", False):
                notes_tf = slide.notes_slide.notes_text_frame
                if notes_tf and notes_tf.text.strip():
                    parts.append(f"(notes) {notes_tf.text.strip()}")
            if parts:
                slide_texts.append(f"[Slide {idx}]\n" + "\n".join(parts))

        text = "\n\n".join(slide_texts)
        chunks = _split_into_chunks(text)
        return ParsedDocument(
            text=text,
            chunks=chunks,
            metadata={
                'filename': filename,
                'filetype': '.pptx',
                'size_bytes': len(content),
                'chunk_count': len(chunks),
                'slide_count': len(slide_texts),
            },
        )


# ---------------------------------------------------------------------------
# HTML parser
# ---------------------------------------------------------------------------


def _html_to_text(raw: str) -> str:
    """Strip an HTML document to readable text (bs4 if available, regex fallback)."""
    try:
        from bs4 import BeautifulSoup  # type: ignore[import]

        soup = BeautifulSoup(raw, 'html.parser')
        for tag in soup(['script', 'style', 'noscript', 'template']):
            tag.decompose()
        return soup.get_text(separator='\n')
    except Exception:
        import html as _html
        import re

        cleaned = re.sub(r'(?is)<(script|style).*?</\1>', ' ', raw)
        cleaned = re.sub(r'(?s)<[^>]+>', ' ', cleaned)
        return _html.unescape(cleaned)


class HtmlParser(BaseParser):
    """Parser for HTML files — strips tags to readable text."""

    supported_extensions = ['.html', '.htm']

    def parse(self, content: bytes, filename: str) -> ParsedDocument:
        try:
            raw = content.decode('utf-8')
        except UnicodeDecodeError:
            raw = content.decode('latin-1', errors='replace')

        text = _html_to_text(raw)
        # Collapse the runs of blank lines bs4/regex stripping tends to produce.
        text = re.sub(r'\n{3,}', '\n\n', text).strip()
        chunks = _split_into_chunks(text)
        return ParsedDocument(
            text=text,
            chunks=chunks,
            metadata={
                'filename': filename,
                'filetype': Path(filename).suffix.lower(),
                'size_bytes': len(content),
                'chunk_count': len(chunks),
            },
        )


# ---------------------------------------------------------------------------
# Best-effort plain-text decode (for unknown / arbitrary file types)
# ---------------------------------------------------------------------------


def best_effort_text(content: bytes, *, max_check: int = 8192) -> str | None:
    """Decode bytes to text if they look textual; return None for binary blobs.

    Used as the fallback for file types with no dedicated parser (source code,
    JSON, YAML, logs, config, …) so "pretty much any text file" indexes without
    a bespoke parser, while true binaries (zip, images, executables) are detected
    and left for the caller to store without indexing.
    """
    if not content:
        return ""

    sample = content[:max_check]
    if b"\x00" in sample:
        return None  # NUL byte → binary

    text: str | None = None
    for enc in ("utf-8", "utf-16", "latin-1"):
        try:
            text = content.decode(enc)
            break
        except (UnicodeDecodeError, LookupError):
            continue
    if text is None:
        return None

    sample_text = text[:max_check]
    if not sample_text:
        return text
    printable = sum(1 for ch in sample_text if ch.isprintable() or ch in "\t\n\r\f\v")
    if printable / len(sample_text) < 0.85:
        return None  # mostly control bytes → treat as binary
    return text


# ---------------------------------------------------------------------------
# Registry
# ---------------------------------------------------------------------------


class DocumentParserRegistry:
    """
    Registry that maps file extensions to parser instances.

    Usage:
        registry = DocumentParserRegistry()
        parser = registry.get_parser("report.txt")
        result = parser.parse(file_bytes, "report.txt")

    Phase 1: .txt, .md, .markdown
    Phase 2: .pdf (pypdf), .docx/.doc (python-docx), .xlsx/.xls/.csv (openpyxl / built-in csv)
    All Phase 2 parsers are registered with graceful degradation — missing deps are skipped.
    """

    def __init__(self):
        self._parsers: dict[str, BaseParser] = {}
        # Register built-in parsers
        self.register(TextParser())

        # Phase 2 — PDF parser (graceful: skipped if pypdf not installed)
        try:
            import pypdf  # noqa: F401
            self.register(PDFParser())
        except ImportError:
            logger.debug("DocumentParserRegistry: pypdf not installed — PDF support disabled")

        # Phase 2 — DOCX parser (graceful: skipped if python-docx not installed)
        try:
            import docx  # noqa: F401
            self.register(DocxParser())
        except ImportError:
            logger.debug("DocumentParserRegistry: python-docx not installed — DOCX support disabled")

        # Phase 2 — XLSX/CSV parser (graceful: openpyxl optional, csv always works)
        xlsx_parser = XlsxParser()
        try:
            import openpyxl  # noqa: F401
            self.register(xlsx_parser)
        except ImportError:
            logger.debug("DocumentParserRegistry: openpyxl not installed — XLSX support disabled, registering CSV only")
            # Register CSV support only (no openpyxl needed for .csv)
            self._parsers['.csv'] = xlsx_parser
            logger.debug("DocumentParserRegistry: registered XlsxParser for '.csv'")

        # Phase 3 — PPTX parser (graceful: skipped if python-pptx not installed)
        try:
            import pptx  # noqa: F401
            self.register(PptxParser())
        except ImportError:
            logger.debug("DocumentParserRegistry: python-pptx not installed — PPTX support disabled")

        # Phase 3 — HTML parser (no hard dependency: bs4 if present, regex fallback otherwise)
        self.register(HtmlParser())

    def register(self, parser: BaseParser) -> None:
        """Register a parser for all its supported extensions."""
        for ext in parser.supported_extensions:
            self._parsers[ext.lower()] = parser
            logger.debug(f"DocumentParserRegistry: registered {type(parser).__name__} for {ext!r}")

    def get_parser(self, filename: str) -> BaseParser:
        """Return the appropriate parser for the given filename.

        Raises UnsupportedFileType if no parser is registered for the extension.
        """
        ext = Path(filename).suffix.lower()
        parser = self._parsers.get(ext)
        if parser is None:
            raise UnsupportedFileType(ext)
        return parser

    @property
    def supported_extensions(self) -> list[str]:
        """List of all supported file extensions."""
        return list(self._parsers.keys())


# ---------------------------------------------------------------------------
# Singleton
# ---------------------------------------------------------------------------

_registry: DocumentParserRegistry | None = None


def get_document_parser_registry() -> DocumentParserRegistry:
    """Return the global DocumentParserRegistry singleton."""
    global _registry
    if _registry is None:
        _registry = DocumentParserRegistry()
    return _registry
